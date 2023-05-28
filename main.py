import collections 
import collections.abc
import os
import json
import requests
import argparse
import tkinter as tk
from tkinter import filedialog
from time import sleep

from pptx import Presentation

parser = argparse.ArgumentParser(
                    prog='.PPTX Fonts parser',
                    description='Get fonts from a presentation and download them')

parser.add_argument('filename', nargs='?')
args = parser.parse_args()

if args.filename is None:
    # если аргументов нет, открыть диалоговое окно выбора файла
    root = tk.Tk()
    root.withdraw()
    file_path = filedialog.askopenfilename()
else:
    # если аргументы есть, использовать переданный файл
    file_path = args.filename


if file_path == "":
    print("[!] Файл не выбран. Выхожу")
    sleep(3)
    exit()
    
print(f'Выбранный файл: {file_path}')

#prs = Presentation(f'{file_path.replace("/", "//")}')
#prs = Presentation(f'{file_path}')
prs = Presentation(file_path)

# создаем пустой список шрифтов
fonts_list = []

# перебираем каждый слайд
for slide in prs.slides:
    # перебираем каждый объект текста на слайде
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # добавляем шрифт, который используется в текством объекте в список
                    font = run.font.name
                    if font not in fonts_list:
                        fonts_list.append(font)

# выводим список всех используемых шрифтов
print("Я нашёл такие шрифты: ")
for i in fonts_list:
    if i != None:
        print(" --- " + i) # Лень убирать None из списка, так бы сделал через "\n".join(fonts_list)


types = ["Bold", "Italic", "Underline", "ExtraBold", 
         "MediumBold", "SemiBold", "SemiboldItalic", 
         "MediumItalic", "ExtraBoldItalic", "RegularBoldItalic", 
         "RegularItalic", "Regular", "Italic", "BoldItalic"]

default_fonts = ["Arial", "Helvetica", "Times New Roman", "Courier New", "Verdana", "Tahoma", "Georgia", "Comic Sans MS"]

# Убираем элементы None и элементы входящие в default_fonts
fonts_list = [font for font in fonts_list if font is not None and font not in default_fonts]

# Вырезаем из элементов типы шрифтов, входящие в types
for i, font in enumerate(fonts_list):
    for f_type in types:
        if f_type in font:
            fonts_list[i] = font.replace(f_type, "").strip()

# Удаляем дублирующиеся элементы
fonts_list = list(set(fonts_list))

print("\nВшух вшух, подчистил шрифты: ")
for i in fonts_list:
    if i != None:
        print(" --- " + i)
print()

if not os.path.exists('temp'):
    os.mkdir('temp')

print("Качаю что могу...")
for font in fonts_list:
    # формируем ссылку на json
    url = f'https://fonts.google.com/download/list?family={font.replace(" ", "%20")}'
   
    # делаем запрос на json
    response = requests.get(url)
    json_file = response.text.replace(")]}'", "")
    # получаем имя файла из json
    fonts_info = json.loads(json_file)
    print(f"'{font}' в прогрессе...")
    for font in fonts_info["manifest"]["fileRefs"]:
        file_name = font['filename'].replace("static/", "")
        
        # получаем ссылку на скачивание из json
        file_url = font['url']
        
        # скачиваем файл
        response = requests.get(file_url)

        # сохраняем файл в папку temp с именем из json
        with open(f'temp/{file_name}', 'wb') as f:
            f.write(response.content)
            
            
print("Загрузка завершена. Выхожу")
sleep(3)